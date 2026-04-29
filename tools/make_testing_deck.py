"""Generate a step-by-step PowerPoint deck from the SpaceboatMania Testing Guide."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DARK_BG = RGBColor(0x14, 0x14, 0x1E)
ACCENT = RGBColor(0x5B, 0x9B, 0xD5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_GRAY = RGBColor(0x99, 0x99, 0x99)
YELLOW = RGBColor(0xFF, 0xD9, 0x66)
GREEN = RGBColor(0x7E, 0xCF, 0x7E)
ORANGE = RGBColor(0xFF, 0xA5, 0x54)
RED_SOFT = RGBColor(0xE8, 0x6B, 0x6B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=None, spacing_pt=6):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        item_color = color
        item_bold = False
        if isinstance(item, tuple):
            item_text, item_color = item[0], item[1]
            if len(item) > 2:
                item_bold = item[2]
        else:
            item_text = item
        p.text = item_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = item_color
        p.font.bold = item_bold
        p.font.name = "Segoe UI"
        p.space_after = Pt(spacing_pt)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buNone = pPr.findall(qn("a:buNone"))
        for bn in buNone:
            pPr.remove(bn)
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "•")
        if bullet_color or item_color != WHITE:
            bc = bullet_color or item_color
            buClr = etree.SubElement(pPr, qn("a:buClr"))
            srgb = etree.SubElement(buClr, qn("a:srgbClr"))
            srgb.set("val", f"{bc.r:02X}{bc.g:02X}{bc.b:02X}" if hasattr(bc, 'r') else str(bc).replace('#',''))
    return tf


def add_phase_badge(slide, phase_num, total=12):
    add_text_box(slide, 10.5, 0.3, 2.5, 0.4,
                 f"Step {phase_num} of {total}", font_size=12, color=DIM_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_exit_criteria(slide, top, items):
    add_text_box(slide, 0.7, top, 5.0, 0.4, "Exit Criteria", font_size=15,
                 color=GREEN, bold=True)
    add_bullet_list(slide, 0.7, top + 0.35, 11.5, 3.0, items, font_size=14,
                    color=GREEN, bullet_color=GREEN)


def add_tip_box(slide, top, text):
    add_text_box(slide, 7.0, top, 5.8, 0.3, "TIP", font_size=12, color=YELLOW, bold=True)
    add_text_box(slide, 7.0, top + 0.25, 5.8, 1.0, text, font_size=13, color=YELLOW)


# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.5, 1.8, 12.3, 1.2, "SpaceboatMania", font_size=52,
             color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 3.2, 12.3, 0.8, "Testing Guide — Step by Step", font_size=30,
             color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 4.3, 12.3, 0.6,
             "Build and test one canonical authored pack end to end", font_size=18,
             color=DIM_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 6.5, 12.3, 0.5,
             "Target route:  main_menu → new game → space → planet POI → atmosphere → "
             "landing region → rooms → combat → MV return → atmosphere → space",
             font_size=13, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 2: Overview / What This Proves ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "What This Guide Proves", font_size=34,
             color=ACCENT, bold=True)
add_bullet_list(slide, 0.7, 1.3, 5.5, 5.0, [
    "Pack creation and pack-aware saves",
    "Authored player, attack, projectile, item, equipment data",
    "Authored UI screens matching the contract",
    "Authored systems, realms, regions, and MV rooms",
    "Authored entities and behaviors",
    "Authored dialogue, shops, and triggers",
    "Save/load and failure recovery paths",
], font_size=18, color=WHITE, bullet_color=ACCENT)

add_text_box(slide, 7.0, 1.3, 5.8, 0.4, "Source of Truth Rules", font_size=20,
             color=ORANGE, bold=True)
add_bullet_list(slide, 7.0, 1.9, 5.8, 4.5, [
    "If editor exposes something not in SUPPORTED_FEATURES.md → bug",
    "If save/pack validation rejects data → blocking",
    "If runtime disagrees with contract after validation → runtime bug",
    "Don't work around validator failures by editing JSON manually",
], font_size=15, color=LIGHT_GRAY, bullet_color=ORANGE)

add_text_box(slide, 7.0, 4.3, 5.8, 0.4, "Prerequisites", font_size=20,
             color=YELLOW, bold=True)
add_bullet_list(slide, 7.0, 4.9, 5.8, 2.0, [
    'Project root: D:/spacegame2/project.godot',
    'Use one fresh pack name (e.g. "golden_path")',
    "Have a Godot 4.x executable available",
], font_size=15, color=LIGHT_GRAY, bullet_color=YELLOW)

# ── Slide 3: Validation Commands ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Validation Commands", font_size=34,
             color=ACCENT, bold=True)
add_text_box(slide, 0.7, 1.1, 12.0, 0.4, "Run these before AND during playtesting — not only at the end",
             font_size=16, color=YELLOW)

cmds = [
    ("Pack validation:", "validate_ui_contract.ps1 -PackId golden_path"),
    ("Smoke-pack validation:", "validate_ui_contract.ps1 -SmokePack"),
    ("Doc drift check:", "validate_ui_contract.ps1 -SyncDocs -CheckDocs"),
    ("Custom Godot path:", "validate_ui_contract.ps1 -GodotBin \"C:\\path\\to\\godot.exe\" -PackId golden_path"),
]
y = 1.7
for label, cmd in cmds:
    add_text_box(slide, 0.7, y, 12.0, 0.3, label, font_size=15, color=ACCENT, bold=True)
    add_text_box(slide, 0.7, y + 0.3, 12.0, 0.4,
                 f"powershell -ExecutionPolicy Bypass -File .\\tools\\{cmd}",
                 font_size=13, color=LIGHT_GRAY, font_name="Cascadia Mono")
    y += 0.85

add_text_box(slide, 0.7, y + 0.2, 12.0, 0.4, "What pack validation checks:", font_size=18,
             color=WHITE, bold=True)
add_bullet_list(slide, 0.7, y + 0.65, 5.5, 3.0, [
    "Authored UI screens & input-map targets",
    "Pack manifest references",
    "System / POI / realm / room references",
    "Room dimensions, doors, hierarchy",
], font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT)
add_bullet_list(slide, 6.5, y + 0.65, 5.5, 3.0, [
    "Trigger, dialogue, shop references",
    "Entity / behavior references",
    "Player-content cross references",
    "Attack, projectile, item, equipment, ability links",
], font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT)

# ── Slide 4: Phase 1 — Create The Pack ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 1)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 1: Create The Pack", font_size=34,
             color=ACCENT, bold=True)
add_text_box(slide, 0.7, 1.2, 5.0, 0.4, "Steps", font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 1.7, 5.5, 2.0, [
    ("From the main menu, click Editor", WHITE),
    ('Create a new pack — e.g. "golden_path"', WHITE),
    ("Confirm the suite shell opens on Campaign tab", WHITE),
], font_size=17, color=WHITE, bullet_color=ACCENT, spacing_pt=10)

add_text_box(slide, 0.7, 3.3, 10.0, 0.4, "Expected Results", font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 3.8, 11.5, 2.5, [
    "Authored data root exists under user://Packs/<pack>/",
    "Top tabs: Campaign, Objects, World, Triggers, UI + FX, Audio, Playtest",
    "Seeded authored UI screens under user://Packs/<pack>/UI/screens/",
    "Default input_map.json present under user://Packs/<pack>/UI/",
], font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT)

add_exit_criteria(slide, 5.2, [
    "Pack directory structure created",
    "All 7 tabs visible in editor suite shell",
    "Seeded UI screens exist on disk",
])

# ── Slide 5: Phase 2 — Author The Player Baseline ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 2)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 2: Author The Player Baseline", font_size=34,
             color=ACCENT, bold=True)
add_text_box(slide, 0.7, 1.1, 5.0, 0.3, "Suite path:  Objects → Player", font_size=16, color=DIM_GRAY)

add_text_box(slide, 0.7, 1.6, 5.5, 0.4, "Minimum Viable Player", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 2.0, 5.5, 2.5, [
    "Pose/sprite data: idle, run, jump, one attack",
    "collision_width",
    "hurtbox_x, hurtbox_y, hurtbox_w, hurtbox_h",
    "weapon_anchor_x, weapon_anchor_y",
], font_size=15, color=LIGHT_GRAY, bullet_color=ACCENT)

add_text_box(slide, 7.0, 1.6, 5.5, 0.4, "Minimum Content Set", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 7.0, 2.0, 5.5, 3.0, [
    ('Attack: "blaster_basic"', LIGHT_GRAY),
    ('Charged attack: "blaster_charged"', LIGHT_GRAY),
    ('Projectile: "bolt_basic"', LIGHT_GRAY),
    ('Item: "medkit"', LIGHT_GRAY),
    ('Equipment: "starter_blaster"', LIGHT_GRAY),
    ('Ability: "double_jump"', LIGHT_GRAY),
], font_size=15, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=8)

add_exit_criteria(slide, 4.8, [
    "Player content saves cleanly",
    "Attack → projectile links validate",
    "Charged attack → attack links validate",
    "Equipment → attack/ability links validate",
    "Item effects within the supported effect surface",
])

# ── Slide 6: Phase 3 — Author The Starting Ship ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 3)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 3: Author The Starting Ship", font_size=34,
             color=ACCENT, bold=True)
add_text_box(slide, 0.7, 1.1, 5.0, 0.3, "Suite path:  Campaign → Starting Ship", font_size=16, color=DIM_GRAY)

add_text_box(slide, 0.7, 1.6, 5.5, 0.4, "Author", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 2.0, 5.5, 2.0, [
    '1 starting ship template ("starter_shuttle")',
    "1 valid core",
    "A minimal module layout",
    "1 starter combat module for opening space segment",
], font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=8)

add_text_box(slide, 7.0, 1.6, 5.5, 0.4, "Also Verify", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 7.0, 2.0, 5.5, 2.5, [
    "Ship builder opens from the suite shell",
    "Ship save persists after close/reopen",
    "Test Fly launches the actual space runtime",
    "Defeat during Test Fly recovers without losing ship state",
], font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=8)

add_exit_criteria(slide, 4.3, [
    "Pack manifest points at a real starting ship template",
    "Ship loads without fallback/default junk",
    "Builder save and reload round-trip correctly",
])

# ── Slide 7: Phase 4 — Author The Minimal UI Slice ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 4)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 4: Author The Minimal UI Slice", font_size=34,
             color=ACCENT, bold=True)
add_text_box(slide, 0.7, 1.1, 5.0, 0.3, "Suite path:  UI + FX → UI + Cinematics", font_size=16, color=DIM_GRAY)

screens_left = [
    ("main_menu", "new_game, load_game, quit_game"),
    ("hud", "HP bar → player.hp/max_hp, weapon label → current_weapon.name, NO buttons"),
    ("pause", "resume, save_game, quit_to_menu"),
    ("inventory", "tab bar, one list/grid, close button"),
]
screens_right = [
    ("shop", "item list → shop.items, buy/sell if needed, close"),
    ("dialogue_box", "speaker label, text area, choice list → choose_dialogue"),
    ("game_over", "load_game or load_slot, quit_to_menu"),
    ("map / boss_intro", "optional but supported"),
]

add_text_box(slide, 0.7, 1.6, 5.5, 0.4, "Required Screens", font_size=18, color=WHITE, bold=True)
y = 2.05
for name, desc in screens_left:
    add_text_box(slide, 0.7, y, 5.5, 0.25, name, font_size=15, color=ACCENT, bold=True)
    add_text_box(slide, 0.7, y + 0.22, 5.5, 0.3, desc, font_size=13, color=LIGHT_GRAY)
    y += 0.55

y = 1.6
add_text_box(slide, 7.0, y, 5.5, 0.4, "More Screens", font_size=18, color=WHITE, bold=True)
y = 2.05
for name, desc in screens_right:
    add_text_box(slide, 7.0, y, 5.5, 0.25, name, font_size=15, color=ACCENT, bold=True)
    add_text_box(slide, 7.0, y + 0.22, 5.5, 0.3, desc, font_size=13, color=LIGHT_GRAY)
    y += 0.55

add_tip_box(slide, 4.5, "For art-driven buttons: use the \"...\" picker on sprite_source,\n"
            "sprite_normal, sprite_hover, sprite_pressed. Use IMPORT in the\n"
            "texture picker for pack-local PNGs. Leave hover/pressed blank to reuse normal.")

add_exit_criteria(slide, 5.6, [
    "Every required screen saves without validation errors",
    "Bindings match SUPPORTED_FEATURES.md contract",
    "No HUD button/item actions authored",
    "Button triggers wired through ui_button / fire_event",
])

# ── Slide 8: Phase 5 — Author The Smallest Useful World ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 5)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 5: Author The World", font_size=34,
             color=ACCENT, bold=True)
add_text_box(slide, 0.7, 1.1, 8.0, 0.3,
             "Suite paths:  World → Systems + Planets  |  World → Realm + Regions",
             font_size=16, color=DIM_GRAY)

add_text_box(slide, 0.7, 1.6, 5.5, 0.4, "Hierarchy", font_size=18, color=WHITE, bold=True)
add_text_box(slide, 0.7, 2.0, 5.5, 0.5,
             "Pack → System → Realm → Region → MV Rooms",
             font_size=20, color=ACCENT, bold=True, font_name="Cascadia Mono")

add_text_box(slide, 0.7, 2.7, 5.5, 0.4, "Author Exactly One Slice", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 3.1, 5.5, 2.0, [
    "1 system, 1 realm, 1 region",
    "4–6 rooms",
    "1 planet POI linking system to realm",
], font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=8)

add_text_box(slide, 7.0, 1.6, 5.5, 0.4, "Recommended Room IDs", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 7.0, 2.0, 5.5, 3.0, [
    "start_room",
    "hallway_a",
    "item_room",
    "shop_room",
    "gate_room",
    "boss_room",
], font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=6)

add_text_box(slide, 7.0, 4.0, 5.5, 0.4, "Setup Notes", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 7.0, 4.4, 5.5, 2.5, [
    "POI realm = your test realm id",
    'Leave "Pack Override" blank (unless testing cross-pack)',
    'Leave "Spawn Room" blank for canonical path',
    "Set valid start_region on the realm",
    "Set atmosphere ground/structure content",
], font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=5)

add_exit_criteria(slide, 6.0, [
    "System / realm / region / room references all validate",
    "Door targets are real rooms  •  POI → realm link validates  •  Region start room resolves",
])

# ── Slide 9: Phase 6 — Environment Authoring ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 6)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 6: Environment Authoring", font_size=34,
             color=ACCENT, bold=True)

add_text_box(slide, 0.7, 1.3, 5.5, 0.4, "For Each Room", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 1.7, 5.5, 2.0, [
    "Paint enough tiles to navigate",
    "Add solid collision",
    "Add one hazard somewhere useful",
    "Set room metadata cleanly",
], font_size=17, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=10)

add_text_box(slide, 7.0, 1.3, 5.5, 0.4, "Recommended Room Roles", font_size=18, color=WHITE, bold=True)
roles = [
    ("start_room", "Safe intro"),
    ("hallway_a", "Traversal"),
    ("item_room", "Progression pickup"),
    ("shop_room", "Dialogue / shop interaction"),
    ("gate_room", "Progression check"),
    ("boss_room", "Combat proof"),
]
y = 1.75
for rid, role in roles:
    add_text_box(slide, 7.0, y, 2.5, 0.3, rid, font_size=14, color=ACCENT, bold=True,
                 font_name="Cascadia Mono")
    add_text_box(slide, 9.6, y, 3.0, 0.3, role, font_size=14, color=LIGHT_GRAY)
    y += 0.35

add_exit_criteria(slide, 4.5, [
    "Room data round-trips through the environment tooling",
    "Collision works  •  Hazards work",
    "Room metadata saves and reloads correctly",
])

# ── Slide 10: Phase 7 — Entity And Behavior ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 7)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 7: Entity And Behavior", font_size=34,
             color=ACCENT, bold=True)

add_text_box(slide, 0.7, 1.3, 5.5, 0.4, "Create", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 1.7, 5.5, 2.0, [
    ('1 enemy definition — "crawler_basic"', LIGHT_GRAY),
    ('1 behavior tree — "crawler_patrol"', LIGHT_GRAY),
], font_size=17, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=10)

add_text_box(slide, 0.7, 2.7, 5.5, 0.4, "Entity Must Define", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 3.1, 5.5, 2.0, [
    "hp",
    "contact_damage",
    "contact_cooldown",
    "move_speed",
], font_size=17, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=6)

add_text_box(slide, 7.0, 1.3, 5.5, 0.4, "Placement", font_size=18, color=WHITE, bold=True)
add_text_box(slide, 7.0, 1.8, 5.5, 0.5,
             "Place the entity in boss_room or a dedicated combat room.",
             font_size=16, color=LIGHT_GRAY)

add_exit_criteria(slide, 4.5, [
    "Entity → behavior reference validates",
    "Room placement references a real entity id",
    "Runtime enemy contact/combat uses authored data, not hardcoded fallback",
])

# ── Slide 11: Phase 8 — Dialogue, Shop, Trigger Chain ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 8)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 8: Dialogue, Shop, Trigger Chain", font_size=34,
             color=ACCENT, bold=True)

add_text_box(slide, 0.7, 1.2, 12.0, 0.4,
             "Build one connected chain — not disconnected samples",
             font_size=18, color=YELLOW, bold=True)

add_text_box(slide, 0.7, 1.8, 5.5, 0.4, "Recommended Chain", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 2.2, 11.0, 3.5, [
    ('1. In shop_room, NPC starts dialogue "shopkeep_intro"', LIGHT_GRAY),
    ("2. Dialogue branches on item / ability / flag / variable condition", LIGHT_GRAY),
    ("3. Dialogue can open a shop or fire a trigger event", LIGHT_GRAY),
    ('4. Shop sells "medkit" and optionally one progression item', LIGHT_GRAY),
    ("5. Trigger in gate_room checks the progression condition", LIGHT_GRAY),
    ("6. One authored return interaction exits MV back to atmosphere overworld", LIGHT_GRAY),
], font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=8)

add_text_box(slide, 0.7, 4.7, 5.5, 0.4, "Authored Return Setup", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 5.1, 11.0, 2.0, [
    "Add an interactable in one MV room",
    'On interact, fire trigger action "return_to_overworld"',
    "Leave region_id, x, y blank for default test path",
], font_size=15, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=6)

add_exit_criteria(slide, 6.0, [
    "Dialogue / shop / trigger validation passes  •  No dangling IDs",
    "Shop item references are real authored items  •  Return path reaches atmosphere overworld",
])

# ── Slide 12: Phase 9 — Save / Load / Failure Recovery ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 9)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 9: Save / Load / Recovery", font_size=34,
             color=ACCENT, bold=True)

add_text_box(slide, 0.7, 1.3, 12.0, 0.4,
             "Add one clear save/checkpoint moment in the route",
             font_size=18, color=WHITE)

add_text_box(slide, 0.7, 2.0, 5.5, 0.4, "At Minimum Prove", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 2.4, 5.5, 2.0, [
    "A save action exists in reachable UI",
    "The authored run can be resumed without manual data edits",
    "Game over can recover through the authored path",
], font_size=17, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=10)

add_text_box(slide, 7.0, 2.0, 5.5, 0.4, "Recommended Checks", font_size=18, color=WHITE, bold=True)
add_bullet_list(slide, 7.0, 2.4, 5.5, 2.0, [
    "Pause-menu save / load",
    "Main-menu load slot path",
    "Game-over load slot path",
], font_size=17, color=LIGHT_GRAY, bullet_color=ACCENT, spacing_pt=10)

add_exit_criteria(slide, 4.3, [
    "Save/load works from reachable authored UI",
    "Game-over recovery is real, not editor-only",
])

# ── Slide 13: Phase 10 — Validate Before Play ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_phase_badge(slide, 10)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Phase 10: Validate Before Play", font_size=34,
             color=ACCENT, bold=True)

add_text_box(slide, 0.7, 1.3, 12.0, 0.4,
             "Before any serious manual play pass:", font_size=18, color=WHITE)
add_bullet_list(slide, 0.7, 1.8, 5.5, 1.5, [
    ("Run pack validation", WHITE, True),
    ("Treat every error as BLOCKING", RED_SOFT, True),
    ("Treat every warning as suspect until explained", YELLOW, True),
], font_size=17, color=WHITE, bullet_color=ACCENT, spacing_pt=10)

add_text_box(slide, 0.7, 3.3, 12.0, 0.4, "Common UI Validation Failures", font_size=18,
             color=WHITE, bold=True)
add_bullet_list(slide, 0.7, 3.7, 5.5, 3.5, [
    "Unknown screen ids in UI/input_map.json",
    "Unsupported bindings",
    "Unsupported actions for a given host",
    "Invalid open_screen targets",
], font_size=15, color=LIGHT_GRAY, bullet_color=RED_SOFT)
add_bullet_list(slide, 6.5, 3.7, 5.5, 3.5, [
    "Missing required action_args",
    "Non-integer args for load_slot or choose_dialogue",
    "Target screen reference that doesn't exist in pack yet",
], font_size=15, color=LIGHT_GRAY, bullet_color=RED_SOFT)

# ── Slide 14: Manual Play Checklist ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Manual Play Checklist", font_size=34,
             color=ACCENT, bold=True)
add_phase_badge(slide, 11)

checklist_left = [
    "1.  Authored main menu appears",
    "2.  new_game starts the authored pack",
    "3.  Authored starting ship loads in space",
    "4.  HUD bindings show live values (no placeholders)",
    "5.  Space flight is playable",
    "6.  Planet POI enters atmosphere overworld",
    "7.  Atmosphere ground content renders",
    "8.  Atmosphere billboards/structures placed correctly",
    "9.  Landing enters region start room in MV",
    "10. Player can move, jump, attack (authored data)",
    "11. Authored attack/projectile behavior occurs",
    "12. Inventory opens and shows authored content",
]
checklist_right = [
    "13. Room transitions work",
    "14. Hazard damage works",
    "15. Enemy contact/combat works",
    "16. Dialogue opens and branches",
    "17. Shop opens and buying works",
    "18. Gate logic respects authored state",
    "19. MV return interaction exits to atmosphere",
    "20. Atmosphere exit returns to space cleanly",
    "21. Pause menu works",
    "22. Save/load works",
    "23. Game-over flow works if triggered",
]

add_bullet_list(slide, 0.5, 1.2, 6.0, 5.5, checklist_left, font_size=14,
                color=LIGHT_GRAY, bullet_color=GREEN, spacing_pt=6)
add_bullet_list(slide, 6.8, 1.2, 6.0, 5.5, checklist_right, font_size=14,
                color=LIGHT_GRAY, bullet_color=GREEN, spacing_pt=6)

# ── Slide 15: Recording Failures + Success Criteria ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.7, 0.4, 10.0, 0.6, "Recording Failures & Success Criteria", font_size=34,
             color=ACCENT, bold=True)
add_phase_badge(slide, 12)

add_text_box(slide, 0.7, 1.2, 5.5, 0.4, "For Every Failure, Record", font_size=18,
             color=RED_SOFT, bold=True)
add_bullet_list(slide, 0.7, 1.6, 5.5, 3.0, [
    "Pack id",
    "Phase of the run",
    "Which editor authored the data",
    "Which file or content id was involved",
    "What you expected vs what happened",
    "Should save validation have caught it?",
    "Should pack validation have caught it?",
], font_size=15, color=LIGHT_GRAY, bullet_color=RED_SOFT, spacing_pt=5)

add_text_box(slide, 0.7, 4.1, 5.5, 0.4, "High-Value Bugs", font_size=18,
             color=ORANGE, bold=True)
add_bullet_list(slide, 0.7, 4.5, 5.5, 1.5, [
    "Runtime failures that should have been rejected at save time",
    "Runtime failures that should have been caught by pack validation",
    "Doc claims that the code does not actually honor",
], font_size=15, color=LIGHT_GRAY, bullet_color=ORANGE, spacing_pt=6)

add_text_box(slide, 7.0, 1.2, 5.5, 0.4, "Success Criteria", font_size=18,
             color=GREEN, bold=True)
add_bullet_list(slide, 7.0, 1.6, 5.5, 3.0, [
    "Canonical pack created entirely through pack-aware tools",
    "Validation passes without unexplained errors",
    "Authored run is playable E2E across space, atmosphere, and MV",
    "UI screens obey the published contract",
    "Docs and runtime behavior match",
    "Remaining failures are narrow bugs, not systemic contract lies",
], font_size=15, color=LIGHT_GRAY, bullet_color=GREEN, spacing_pt=6)

add_text_box(slide, 7.0, 4.1, 5.5, 0.4, "What NOT To Do Yet", font_size=18,
             color=RED_SOFT, bold=True)
add_bullet_list(slide, 7.0, 4.5, 5.5, 2.5, [
    "Build a large world",
    "Polish visuals heavily",
    "Create many enemies before one path is proven",
    "Create many items/abilities before one progression path works",
    "Rely on unsupported screen ids, bindings, or actions",
], font_size=15, color=LIGHT_GRAY, bullet_color=RED_SOFT, spacing_pt=5)

# ── Save ──
out_path = r"D:\spacegame2\SpaceboatMania_Testing_Guide.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
